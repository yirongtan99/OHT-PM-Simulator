from pptx import Presentation
import sys

def read_pptx(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.replace("\n", " ").strip())
                    print(" | ".join(row_data))

if __name__ == "__main__":
    read_pptx(sys.argv[1])
